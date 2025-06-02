from crewai.tools import BaseTool
from typing import Optional, Type, Tuple, List, Dict, Any, Union
from pydantic import BaseModel, Field, ValidationError
import logging
import os
from pathlib import Path
import uuid
import traceback
from datetime import datetime
from io import BytesIO
import json

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, PP_PARAGRAPH_ALIGNMENT, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR_TYPE
from pptx.enum.dml import MSO_FILL, MSO_THEME_COLOR, MSO_LINE, MSO_LINE_DASH_STYLE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_LABEL_POSITION, XL_DATA_LABEL_POSITION
from pptx.enum.action import PP_ACTION
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
from pptx.shapes.connector import Connector
from pptx.shapes.freeform import FreeformBuilder

# Define the output schema class directly to avoid import issues
class PythonPPTXToolOutput(BaseModel):
    """Output schema for PythonPPTXTool."""
    success: bool = Field(description="Whether the operation was successful")
    message: str = Field(description="Status message")
    file_path: str = Field(description="Absolute path to the created presentation file")
    relative_path: str = Field(description="Relative path to the created presentation file")
    content: str = Field(description="Content used to create the presentation")
    title: str = Field(description="Title of the presentation")

# Configure logger
logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# Define Pydantic models for structured input
class Headline(BaseModel):
    title: str
    subtitle: Optional[str] = None
    date: Union[datetime, str] = Field(default_factory=datetime.now)

    # For JSON serialization/deserialization
    @classmethod
    def parse_date(cls, date_str):
        if date_str:
            try:
                return datetime.fromisoformat(date_str.replace('Z', '+00:00'))
            except ValueError:
                for fmt in ['%Y-%m-%d', '%d/%m/%Y', '%m/%d/%Y', '%B %d, %Y']:
                    try:
                        return datetime.strptime(date_str, fmt)
                    except ValueError:
                        continue
        return datetime.now()


class TextFormatting(BaseModel):
    bold: bool = False
    italic: bool = False
    underline: bool = False
    font_size: Optional[int] = None
    font_name: Optional[str] = None
    color: Optional[Union[str, Tuple[int, int, int]]] = None
    alignment: Optional[str] = None  # "left", "center", "right", "justify"


class BulletPoint(BaseModel):
    text: str
    level: int = 0
    bold: bool = False
    italic: bool = False
    underline: bool = False
    font_size: Optional[int] = None
    font_name: Optional[str] = None
    color: Optional[Union[str, Tuple[int, int, int]]] = None


class TextRun(BaseModel):
    text: str
    formatting: Optional[TextFormatting] = None


class Paragraph(BaseModel):
    text: Optional[str] = None
    level: int = 0
    alignment: Optional[str] = None  # "left", "center", "right", "justify"
    formatting: Optional[TextFormatting] = None
    runs: Optional[List[TextRun]] = None  # For mixed formatting


class TextFrame(BaseModel):
    paragraphs: List[Paragraph] = []
    auto_size: Optional[bool] = True
    word_wrap: bool = True
    vertical_alignment: Optional[str] = None  # "top", "middle", "bottom"
    margin_left: Optional[float] = None  # In inches
    margin_right: Optional[float] = None
    margin_top: Optional[float] = None
    margin_bottom: Optional[float] = None


class AxisProperties(BaseModel):
    title: Optional[str] = None
    min_value: Optional[float] = None
    max_value: Optional[float] = None
    major_unit: Optional[float] = None
    minor_unit: Optional[float] = None
    has_major_gridlines: bool = False
    has_minor_gridlines: bool = False
    tick_label_position: Optional[str] = None  # "high", "low", "next_to", "none"
    number_format: Optional[str] = None  # "0", "0.00", "0%", etc.


class DataLabels(BaseModel):
    enabled: bool = False
    position: Optional[str] = None  # "inside_end", "outside_end", "center", etc.
    show_values: bool = True
    show_percentage: bool = False
    show_category_name: bool = False
    show_series_name: bool = False
    number_format: Optional[str] = None


class ChartSeries(BaseModel):
    name: str
    values: List[float]
    line_properties: Optional[Dict[str, Any]] = None  # For line charts
    marker_properties: Optional[Dict[str, Any]] = None  # For line charts
    fill_color: Optional[Union[str, Tuple[int, int, int]]] = None
    data_labels: Optional[DataLabels] = None


class ChartData(BaseModel):
    title: str
    chart_type: str = "BAR"  # BAR, LINE, PIE, AREA, XY, BUBBLE, etc.
    categories: Optional[List[str]] = None  # Not needed for XY/Bubble
    series: Optional[Dict[str, List[float]]] = None  # For category charts
    xy_series: Optional[List[Dict[str, Any]]] = None  # For XY charts
    bubble_series: Optional[List[Dict[str, Any]]] = None  # For bubble charts
    has_legend: bool = True
    legend_position: Optional[str] = "RIGHT"  # RIGHT, LEFT, TOP, BOTTOM
    axis_properties: Optional[Dict[str, AxisProperties]] = None
    data_labels: Optional[DataLabels] = None
    has_data_table: bool = False
    style: Optional[int] = None  # 1-48


class TableData(BaseModel):
    headers: List[str]
    rows: List[List[str]]
    style: Optional[str] = "Medium Style 2 Accent 1"  # Default PowerPoint table style
    first_row_header: bool = True
    cell_formatting: Optional[List[Dict[str, Any]]] = None  # For specific cell formatting
    column_widths: Optional[List[float]] = None  # In inches


class ShapeProperties(BaseModel):
    fill_type: Optional[str] = "solid"  # solid, gradient, picture, pattern, none
    fill_color: Optional[Union[str, Tuple[int, int, int]]] = None
    gradient_stops: Optional[List[Tuple[float, Tuple[int, int, int]]]] = None
    line_color: Optional[Union[str, Tuple[int, int, int]]] = None
    line_width: Optional[float] = None  # In points
    line_dash: Optional[str] = None  # solid, dash, dot, etc.
    line_style: Optional[str] = None
    shadow: Optional[bool] = False
    shadow_properties: Optional[Dict[str, Any]] = None
    rotation: Optional[float] = None  # In degrees
    transparency: Optional[float] = None  # 0-1


class GroupShapeItem(BaseModel):
    type: str  # "auto_shape", "picture", "chart", "table", "text_box"
    shape_data: Dict[str, Any]
    position_x: float  # In inches
    position_y: float  # In inches
    width: float  # In inches
    height: float  # In inches
    rotation: Optional[float] = None  # In degrees
    shape_properties: Optional[ShapeProperties] = None
    text_frame: Optional[TextFrame] = None


class HyperlinkAction(BaseModel):
    url: str
    tooltip: Optional[str] = None
    target: Optional[str] = None  # "_blank", "_self", etc.


class SlideAction(BaseModel):
    action_type: str = "hyperlink"  # "hyperlink", "slide", "program", "macro"
    action_data: Any  # URL for hyperlink, slide index for slide action, etc.
    tooltip: Optional[str] = None


class PicturePlaceholder(BaseModel):
    image_path: str
    alt_text: Optional[str] = None
    crop: Optional[Dict[str, float]] = None
    brightness: Optional[float] = None  # -1.0 to 1.0
    contrast: Optional[float] = None  # -1.0 to 1.0


class ConnectorShape(BaseModel):
    connector_type: str = "straight"  # straight, elbow, curved
    begin_point: Tuple[float, float]  # x, y in inches
    end_point: Tuple[float, float]  # x, y in inches
    begin_connect: Optional[Dict[str, Any]] = None  # Shape ID and connection point
    end_connect: Optional[Dict[str, Any]] = None  # Shape ID and connection point
    line_properties: Optional[Dict[str, Any]] = None


class ContentSlide(BaseModel):
    title: str
    subtitle: Optional[str] = None
    bullet_points: Optional[List[Union[str, Dict, BulletPoint]]] = None
    bullets: Optional[List[Union[str, Dict, BulletPoint]]] = None  # For backward compatibility
    paragraphs: Optional[List[Paragraph]] = None  # For rich text
    image_path: Optional[str] = None
    chart_data: Optional[Union[Dict, ChartData]] = None
    table_data: Optional[Union[Dict, TableData]] = None
    shapes: Optional[List[Union[Dict, ShapeProperties]]] = None
    group_shapes: Optional[List[GroupShapeItem]] = None
    connectors: Optional[List[ConnectorShape]] = None
    notes: Optional[str] = None  # Speaker notes
    content: Optional[str] = None  # Simple content field
    background: Optional[ShapeProperties] = None
    hyperlink: Optional[HyperlinkAction] = None
    slide_action: Optional[SlideAction] = None


class CoreProperties(BaseModel):
    title: Optional[str] = None
    author: Optional[str] = None
    subject: Optional[str] = None
    category: Optional[str] = None
    comments: Optional[str] = None
    keywords: Optional[Union[str, List[str]]] = None
    language: Optional[str] = None
    revision: Optional[int] = None
    status: Optional[str] = None
    version: Optional[str] = None
    last_modified_by: Optional[str] = None
    created: Optional[Union[datetime, str]] = None
    modified: Optional[Union[datetime, str]] = None


class PresentationContent(BaseModel):
    title: str
    author: Optional[str] = None
    headline: Optional[Dict] = None
    slides: List[Dict]
    theme_color: Optional[str] = "blue"
    include_footer: bool = True
    revision: int = 1
    company: Optional[str] = None
    keywords: Optional[List[str]] = None
    core_properties: Optional[CoreProperties] = None
    slide_master: Optional[Dict] = None
    slide_layouts: Optional[List[Dict]] = None


class PythonPPTXInput(BaseModel):
    """Input schema for PythonPPTXTool."""
    content: str = Field(..., description="The raw content to be converted into slides. Can be plain text or JSON format.")
    title: str = Field(default="New Presentation", description="The title of the presentation")
    output_dir: str = Field(default="./presentations", description="Directory to save the presentation (Note: This will be ignored and ./output will be used instead)")
    template_path: Optional[str] = Field(default=None, description="Path to a PPTX template file (optional)")


def configure_text_frame(text_frame, auto_size=True, word_wrap=True, vertical_alignment=None, 
                       margins=None):
    """Configure text frame properties for better content display.
    
    Args:
        text_frame: The text frame to configure
        auto_size: Whether to enable auto-sizing of text
        word_wrap: Whether to enable word wrapping
        vertical_alignment: Vertical alignment of text (optional)
        margins: Dictionary with margin values for left, right, top, bottom in inches
    """
    if auto_size:
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    
    text_frame.word_wrap = word_wrap
    
    # Set default margins
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)
    text_frame.margin_top = Inches(0.05)
    text_frame.margin_bottom = Inches(0.05)
    
    # Apply custom margins if provided
    if margins:
        if 'left' in margins:
            text_frame.margin_left = Inches(margins['left'])
        if 'right' in margins:
            text_frame.margin_right = Inches(margins['right'])
        if 'top' in margins:
            text_frame.margin_top = Inches(margins['top'])
        if 'bottom' in margins:
            text_frame.margin_bottom = Inches(margins['bottom'])
    
    # Set vertical alignment if provided
    if vertical_alignment:
        if vertical_alignment.lower() == 'top':
            text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        elif vertical_alignment.lower() == 'middle':
            text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        elif vertical_alignment.lower() == 'bottom':
            text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.BOTTOM


def add_formatted_text(text_frame, text, alignment=PP_ALIGN.LEFT, level=0, 
                      bold=False, italic=False, underline=False, color=None, size=None, 
                      font_name=None):
    """Add a paragraph with formatting to a text frame.
    
    Args:
        text_frame: The text frame to add text to
        text: The text content
        alignment: Text alignment (LEFT, CENTER, RIGHT, JUSTIFY)
        level: Indentation level for bullets
        bold: Whether text should be bold
        italic: Whether text should be italic
        underline: Whether text should be underlined
        color: RGB color tuple or RGBColor object
        size: Font size in points
        font_name: Font name
    
    Returns:
        The newly added paragraph
    """
    p = text_frame.add_paragraph()
    p.text = text
    p.alignment = alignment
    p.level = level
    
    # Apply font formatting
    p.font.bold = bold
    p.font.italic = italic
    if underline:
        p.font.underline = True
    
    if color:
        if isinstance(color, tuple) and len(color) == 3:
            p.font.color.rgb = RGBColor(*color)
        elif isinstance(color, RGBColor):
            p.font.color.rgb = color
        elif isinstance(color, str) and color.startswith('#'):
            # Convert hex color to RGB
            r = int(color[1:3], 16)
            g = int(color[3:5], 16)
            b = int(color[5:7], 16)
            p.font.color.rgb = RGBColor(r, g, b)
    
    if size is not None:
        p.font.size = Pt(size)
    
    if font_name is not None:
        p.font.name = font_name
    
    return p


def add_text_with_mixed_formatting(text_frame, paragraph_data):
    """Add text with mixed formatting within a single paragraph.
    
    Args:
        text_frame: The text frame to add text to
        paragraph_data: Dictionary or Paragraph object with text and formatting
        
    Returns:
        The newly added paragraph
    """
    if isinstance(paragraph_data, Paragraph):
        # Use pydantic model
        p_model = paragraph_data
    else:
        # Create from dict
        p_model = Paragraph(**paragraph_data)
    
    # Create paragraph and apply paragraph-level formatting
    p = text_frame.add_paragraph()
    p.alignment = get_pp_alignment(p_model.alignment) if p_model.alignment else PP_ALIGN.LEFT
    p.level = p_model.level
    
    # If we have runs with different formatting
    if p_model.runs and len(p_model.runs) > 0:
        # Clear default text
        p.text = ""
        
        # Add each run with its specific formatting
        for run_data in p_model.runs:
            run = p.add_run()
            run.text = run_data.text
            
            # Apply formatting if specified
            if run_data.formatting:
                fmt = run_data.formatting
                if fmt.bold is not None:
                    run.font.bold = fmt.bold
                if fmt.italic is not None:
                    run.font.italic = fmt.italic
                if fmt.underline is not None:
                    run.font.underline = fmt.underline
                if fmt.font_size is not None:
                    run.font.size = Pt(fmt.font_size)
                if fmt.font_name is not None:
                    run.font.name = fmt.font_name
                if fmt.color is not None:
                    set_font_color(run.font, fmt.color)
    else:
        # Just set the text and apply formatting to the whole paragraph
        p.text = p_model.text or ""
        
        # Apply paragraph-level formatting
        if p_model.formatting:
            fmt = p_model.formatting
            if fmt.bold is not None:
                p.font.bold = fmt.bold
            if fmt.italic is not None:
                p.font.italic = fmt.italic
            if fmt.underline is not None:
                p.font.underline = fmt.underline
            if fmt.font_size is not None:
                p.font.size = Pt(fmt.font_size)
            if fmt.font_name is not None:
                p.font.name = fmt.font_name
            if fmt.color is not None:
                set_font_color(p.font, fmt.color)
    
    return p


def set_font_color(font, color):
    """Helper function to set font color from various formats.
    
    Args:
        font: The font object to modify
        color: Color as RGB tuple, RGBColor object, or hex string
    """
    if isinstance(color, tuple) and len(color) == 3:
        font.color.rgb = RGBColor(*color)
    elif isinstance(color, RGBColor):
        font.color.rgb = color
    elif isinstance(color, str) and color.startswith('#'):
        # Convert hex color to RGB
        r = int(color[1:3], 16)
        g = int(color[3:5], 16)
        b = int(color[5:7], 16)
        font.color.rgb = RGBColor(r, g, b)
    elif isinstance(color, str) and hasattr(MSO_THEME_COLOR, color.upper()):
        # Use theme color
        font.color.theme_color = getattr(MSO_THEME_COLOR, color.upper())


def get_pp_alignment(alignment_str):
    """Convert string alignment to PowerPoint enum value.
    
    Args:
        alignment_str: String representation of alignment
        
    Returns:
        PP_ALIGN enumeration value
    """
    if not alignment_str:
        return PP_ALIGN.LEFT
        
    mapping = {
        'left': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT,
        'justify': PP_ALIGN.JUSTIFY
    }
    
    return mapping.get(alignment_str.lower(), PP_ALIGN.LEFT)


def add_connector(slide, connector_type, begin_point, end_point, 
                 line_color=None, line_width=None, line_dash=None):
    """Add a connector shape to connect elements.
    
    Args:
        slide: The slide to add the connector to
        connector_type: Type of connector (straight, elbow, curved)
        begin_point: Starting point as (x, y) tuple in inches
        end_point: Ending point as (x, y) tuple in inches
        line_color: Color for the connector line
        line_width: Width of the connector line in points
        line_dash: Dash style for the line
        
    Returns:
        The created connector shape
    """
    # Map connector type to MSO enum
    connector_type_map = {
        'straight': MSO_CONNECTOR_TYPE.STRAIGHT,
        'elbow': MSO_CONNECTOR_TYPE.ELBOW,
        'curved': MSO_CONNECTOR_TYPE.CURVE
    }
    
    mso_connector_type = connector_type_map.get(
        connector_type.lower(), MSO_CONNECTOR_TYPE.STRAIGHT
    )
    
    # Create connector
    connector = slide.shapes.add_connector(
        mso_connector_type,
        Inches(begin_point[0]), Inches(begin_point[1]),
        Inches(end_point[0]), Inches(end_point[1])
    )
    
    # Apply line formatting
    if line_color:
        set_shape_line_color(connector, line_color)
    
    if line_width is not None:
        connector.line.width = Pt(line_width)
    
    if line_dash:
        set_line_dash_style(connector.line, line_dash)
    
    return connector


def set_shape_line_color(shape, color):
    """Helper function to set shape line color from various formats.
    
    Args:
        shape: The shape object to modify
        color: Color as RGB tuple, RGBColor object, or hex string
    """
    if isinstance(color, tuple) and len(color) == 3:
        shape.line.color.rgb = RGBColor(*color)
    elif isinstance(color, RGBColor):
        shape.line.color.rgb = color
    elif isinstance(color, str) and color.startswith('#'):
        # Convert hex color to RGB
        r = int(color[1:3], 16)
        g = int(color[3:5], 16)
        b = int(color[5:7], 16)
        shape.line.color.rgb = RGBColor(r, g, b)
    elif isinstance(color, str) and hasattr(MSO_THEME_COLOR, color.upper()):
        # Use theme color
        shape.line.color.theme_color = getattr(MSO_THEME_COLOR, color.upper())


def set_line_dash_style(line, dash_style):
    """Set line dash style based on string description.
    
    Args:
        line: The line object to modify
        dash_style: String description of dash style
    """
    dash_mapping = {
        'solid': MSO_LINE_DASH_STYLE.SOLID,
        'round_dot': MSO_LINE_DASH_STYLE.ROUND_DOT,
        'square_dot': MSO_LINE_DASH_STYLE.SQUARE_DOT,
        'dash': MSO_LINE_DASH_STYLE.DASH,
        'dash_dot': MSO_LINE_DASH_STYLE.DASH_DOT,
        'long_dash': MSO_LINE_DASH_STYLE.LONG_DASH,
        'long_dash_dot': MSO_LINE_DASH_STYLE.LONG_DASH_DOT,
        'long_dash_dot_dot': MSO_LINE_DASH_STYLE.LONG_DASH_DOT_DOT,
    }
    
    if dash_style.lower() in dash_mapping:
        line.dash_style = dash_mapping[dash_style.lower()]
    

def create_freeform_shape(slide, vertices, fill_color=None, line_color=None, line_width=None):
    """Create a freeform shape with custom geometry.
    
    Args:
        slide: The slide to add the shape to
        vertices: List of (x, y) coordinates in inches defining the shape
        fill_color: Fill color for the shape
        line_color: Line color for the shape outline
        line_width: Line width for the shape outline in points
        
    Returns:
        The created freeform shape
    """
    if not vertices or len(vertices) < 3:
        raise ValueError("At least 3 vertices are required to create a freeform shape")
    
    # Create a freeform builder
    builder = FreeformBuilder(slide.shapes)
    
    # Set starting point
    start_x, start_y = vertices[0]
    builder.start_free_form(Inches(start_x), Inches(start_y))
    
    # Add line segments
    for x, y in vertices[1:]:
        builder.add_line_segment(Inches(x), Inches(y))
    
    # Close the shape
    shape = builder.convert_to_shape()
    
    # Apply formatting
    if fill_color:
        set_shape_fill_color(shape, fill_color)
    
    if line_color:
        set_shape_line_color(shape, line_color)
    
    if line_width is not None:
        shape.line.width = Pt(line_width)
    
    return shape


def set_shape_fill_color(shape, color):
    """Helper function to set shape fill color from various formats.
    
    Args:
        shape: The shape object to modify
        color: Color as RGB tuple, RGBColor object, or hex string
    """
    shape.fill.solid()
    
    if isinstance(color, tuple) and len(color) == 3:
        shape.fill.fore_color.rgb = RGBColor(*color)
    elif isinstance(color, RGBColor):
        shape.fill.fore_color.rgb = color
    elif isinstance(color, str) and color.startswith('#'):
        # Convert hex color to RGB
        r = int(color[1:3], 16)
        g = int(color[3:5], 16)
        b = int(color[5:7], 16)
        shape.fill.fore_color.rgb = RGBColor(r, g, b)
    elif isinstance(color, str) and hasattr(MSO_THEME_COLOR, color.upper()):
        # Use theme color
        shape.fill.fore_color.theme_color = getattr(MSO_THEME_COLOR, color.upper())


def apply_gradient_fill(shape, start_color, end_color, gradient_type='linear', angle=90):
    """Apply gradient fill to a shape.
    
    Args:
        shape: The shape to apply gradient to
        start_color: Start color for gradient (RGB tuple or hex string)
        end_color: End color for gradient (RGB tuple or hex string)
        gradient_type: Type of gradient ('linear' or 'radial')
        angle: Angle for linear gradient in degrees
        
    Returns:
        The modified shape
    """
    shape.fill.gradient()
    
    # Set start color
    if isinstance(start_color, tuple) and len(start_color) == 3:
        shape.fill.gradient_stops[0].color.rgb = RGBColor(*start_color)
    elif isinstance(start_color, str) and start_color.startswith('#'):
        r = int(start_color[1:3], 16)
        g = int(start_color[3:5], 16)
        b = int(start_color[5:7], 16)
        shape.fill.gradient_stops[0].color.rgb = RGBColor(r, g, b)
    
    # Set end color
    if isinstance(end_color, tuple) and len(end_color) == 3:
        shape.fill.gradient_stops[1].color.rgb = RGBColor(*end_color)
    elif isinstance(end_color, str) and end_color.startswith('#'):
        r = int(end_color[1:3], 16)
        g = int(end_color[3:5], 16)
        b = int(end_color[5:7], 16)
        shape.fill.gradient_stops[1].color.rgb = RGBColor(r, g, b)
    
    # Set angle for linear gradient
    if gradient_type.lower() == 'linear':
        shape.fill.gradient_angle = angle
    
    return shape


def add_hyperlink(shape, url, tooltip=None, target=None):
    """Add a hyperlink to a shape.
    
    Args:
        shape: The shape to add the hyperlink to
        url: The URL for the hyperlink
        tooltip: Optional tooltip/screen tip text
        target: Target for the hyperlink (not directly supported by python-pptx)
        
    Returns:
        The modified shape
    """
    # Configure the click action
    click_action = shape.click_action
    click_action.action = PP_ACTION.HYPERLINK
    click_action.hyperlink.address = url
    
    # Set tooltip if provided
    if tooltip:
        click_action.hyperlink.screen_tip = tooltip
    
    return shape


def create_xy_chart(slide, x, y, width, height, chart_data, title=None, 
                   has_legend=True, legend_position=None, style=None):
    """Create an XY (scatter) chart.
    
    Args:
        slide: The slide to add the chart to
        x, y: Position (in inches)
        width, height: Size (in inches)
        chart_data: List of dictionaries containing series data
        title: Title for the chart
        has_legend: Whether to display a legend
        legend_position: Position of the legend
        style: Chart style (1-48)
        
    Returns:
        The created chart
    """
    # Create chart data object for XY chart
    cd = XyChartData()
    
    # Add series data
    for series_dict in chart_data:
        series_name = series_dict.get('name', 'Series')
        series = cd.add_series(series_name)
        
        # Add data points
        if 'points' in series_dict:
            for point in series_dict['points']:
                if isinstance(point, tuple) and len(point) >= 2:
                    series.add_data_point(point[0], point[1])
                elif isinstance(point, dict) and 'x' in point and 'y' in point:
                    series.add_data_point(point['x'], point['y'])
    
    # Create chart on slide
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER, 
        Inches(x), Inches(y), 
        Inches(width), Inches(height), 
        cd
    ).chart
    
    # Set chart title
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
    
    # Configure legend
    chart.has_legend = has_legend
    if has_legend and legend_position:
        legend_pos_map = {
            'right': XL_LEGEND_POSITION.RIGHT,
            'left': XL_LEGEND_POSITION.LEFT,
            'top': XL_LEGEND_POSITION.TOP,
            'bottom': XL_LEGEND_POSITION.BOTTOM,
            'corner': XL_LEGEND_POSITION.CORNER
        }
        if isinstance(legend_position, str) and legend_position.lower() in legend_pos_map:
            chart.legend.position = legend_pos_map[legend_position.lower()]
    
    # Set chart style if provided
    if style is not None:
        chart.style = style
    
    return chart


def create_bubble_chart(slide, x, y, width, height, chart_data, title=None, 
                        has_legend=True, legend_position=None, style=None):
    """Create a bubble chart.
    
    Args:
        slide: The slide to add the chart to
        x, y: Position (in inches)
        width, height: Size (in inches)
        chart_data: List of dictionaries containing series data
        title: Title for the chart
        has_legend: Whether to display a legend
        legend_position: Position of the legend
        style: Chart style (1-48)
        
    Returns:
        The created chart
    """
    # Create chart data object for bubble chart
    cd = BubbleChartData()
    
    # Add series data
    for series_dict in chart_data:
        series_name = series_dict.get('name', 'Series')
        series = cd.add_series(series_name)
        
        # Add data points
        if 'points' in series_dict:
            for point in series_dict['points']:
                if isinstance(point, tuple) and len(point) >= 3:
                    # (x, y, size)
                    series.add_data_point(point[0], point[1], point[2])
                elif isinstance(point, dict) and 'x' in point and 'y' in point and 'size' in point:
                    series.add_data_point(point['x'], point['y'], point['size'])
    
    # Create chart on slide
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BUBBLE, 
        Inches(x), Inches(y), 
        Inches(width), Inches(height), 
        cd
    ).chart
    
    # Set chart title
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
    
    # Configure legend
    chart.has_legend = has_legend
    if has_legend and legend_position:
        legend_pos_map = {
            'right': XL_LEGEND_POSITION.RIGHT,
            'left': XL_LEGEND_POSITION.LEFT,
            'top': XL_LEGEND_POSITION.TOP,
            'bottom': XL_LEGEND_POSITION.BOTTOM,
            'corner': XL_LEGEND_POSITION.CORNER
        }
        if isinstance(legend_position, str) and legend_position.lower() in legend_pos_map:
            chart.legend.position = legend_pos_map[legend_position.lower()]
    
    # Set chart style if provided
    if style is not None:
        chart.style = style
    
    return chart


def customize_chart_axis(chart, axis_type, properties):
    """Customize a chart axis with detailed properties.
    
    Args:
        chart: The chart to modify
        axis_type: 'category' or 'value'
        properties: Dictionary of axis properties
        
    Returns:
        The modified axis
    """
    # Get the appropriate axis
    if axis_type.lower() == 'category':
        axis = chart.category_axis
    else:  # value
        axis = chart.value_axis
    
    # Set title if provided
    if 'title' in properties:
        axis.has_title = True
        axis.axis_title.text_frame.text = properties['title']
    
    # Set axis scale for value axis
    if axis_type.lower() == 'value':
        if 'min_value' in properties:
            axis.minimum_scale = properties['min_value']
        
        if 'max_value' in properties:
            axis.maximum_scale = properties['max_value']
        
        if 'major_unit' in properties:
            axis.major_unit = properties['major_unit']
        
        if 'minor_unit' in properties:
            axis.minor_unit = properties['minor_unit']
    
    # Set gridlines
    if 'has_major_gridlines' in properties:
        axis.has_major_gridlines = properties['has_major_gridlines']
    
    if 'has_minor_gridlines' in properties:
        axis.has_minor_gridlines = properties['has_minor_gridlines']
    
    # Set tick label position
    if 'tick_label_position' in properties:
        pos = properties['tick_label_position'].upper()
        if hasattr(XL_TICK_LABEL_POSITION, pos):
            axis.tick_label_position = getattr(XL_TICK_LABEL_POSITION, pos)
    
    # Set number format
    if 'number_format' in properties and hasattr(axis.tick_labels, 'number_format'):
        axis.tick_labels.number_format = properties['number_format']
    
    return axis


def configure_data_labels(chart, data_label_properties):
    """Configure data labels for a chart.
    
    Args:
        chart: The chart to modify
        data_label_properties: Dictionary of data label properties
        
    Returns:
        The modified chart
    """
    # Check if we have a plot to work with
    if not hasattr(chart, 'plots') or not chart.plots:
        return chart
    
    plot = chart.plots[0]
    plot.has_data_labels = data_label_properties.get('enabled', False)
    
    if plot.has_data_labels:
        data_labels = plot.data_labels
        
        # Set position if specified
        position = data_label_properties.get('position')
        if position:
            pos_map = {
                'center': XL_DATA_LABEL_POSITION.CENTER,
                'inside_end': XL_DATA_LABEL_POSITION.INSIDE_END,
                'inside_base': XL_DATA_LABEL_POSITION.INSIDE_BASE,
                'outside_end': XL_DATA_LABEL_POSITION.OUTSIDE_END,
                'above': XL_DATA_LABEL_POSITION.ABOVE,
                'below': XL_DATA_LABEL_POSITION.BELOW,
                'best_fit': XL_DATA_LABEL_POSITION.BEST_FIT,
                'left': XL_DATA_LABEL_POSITION.LEFT,
                'right': XL_DATA_LABEL_POSITION.RIGHT,
            }
            if position.lower() in pos_map:
                data_labels.position = pos_map[position.lower()]
        
        # Configure what to show
        if 'show_values' in data_label_properties:
            data_labels.show_value = data_label_properties['show_values']
        
        if 'show_percentage' in data_label_properties:
            data_labels.show_percentage = data_label_properties['show_percentage']
        
        if 'show_category_name' in data_label_properties:
            data_labels.show_category_name = data_label_properties['show_category_name']
        
        if 'show_series_name' in data_label_properties:
            data_labels.show_series_name = data_label_properties['show_series_name']
        
        # Set number format if specified
        if 'number_format' in data_label_properties:
            data_labels.number_format = data_label_properties['number_format']
    
    return chart


def create_title_slide(prs, content):
    """Create the title slide with headline"""
    # Check if headline exists in content
    if not hasattr(content, 'headline') or content.headline is None:
        # Create a simple title slide without headline data
        try:
            title_slide_layout = prs.slide_layouts[0]  # Title Slide layout
            slide = prs.slides.add_slide(title_slide_layout)
            
            # Set title if title placeholder exists
            try:
                title = slide.shapes.title
                title.text = content.title
                configure_text_frame(title.text_frame)
                title.text_frame.paragraphs[0].font.size = Pt(44)
                title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 112, 192)  # Blue color
            except (AttributeError, IndexError):
                # Create a textbox for title if no title placeholder
                left = Inches(1)
                top = Inches(1)
                width = Inches(8)
                height = Inches(1.2)
                
                title_box = slide.shapes.add_textbox(left, top, width, height)
                text_frame = title_box.text_frame
                configure_text_frame(text_frame)
                text_frame.text = content.title
                text_frame.paragraphs[0].font.size = Pt(44)
                text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 112, 192)
                text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        except Exception as e:
            # Fall back to blank slide if all else fails
            logger.warning(f"Warning: {e}. Using blank slide for title.")
            slide = prs.slides.add_slide(prs.slide_layouts[-1])  # Usually blank layout
            
            # Add title textbox
            left = Inches(1)
            top = Inches(1)
            width = Inches(8)
            height = Inches(1.2)
            
            title_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = title_box.text_frame
            configure_text_frame(text_frame)
            text_frame.text = content.title
            text_frame.paragraphs[0].font.size = Pt(44)
            text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 112, 192)
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        return slide
        
    # Try to use title slide layout
    try:
        title_slide_layout = prs.slide_layouts[0]  # Title Slide layout
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Extract headline data once
        headline = Headline(**content.headline)
        
        # Set title if title placeholder exists
        try:
            title = slide.shapes.title
            title.text = headline.title
            configure_text_frame(title.text_frame)
            title.text_frame.paragraphs[0].font.size = Pt(44)
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 112, 192)  # Blue color
        except (AttributeError, IndexError):
            # Create a textbox for title if no title placeholder
            left = Inches(1)
            top = Inches(1)
            width = Inches(8)
            height = Inches(1.2)
            
            title_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = title_box.text_frame
            configure_text_frame(text_frame)
            text_frame.text = headline.title
            text_frame.paragraphs[0].font.size = Pt(44)
            text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 112, 192)
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Try to use subtitle placeholder
        try:
            # First, get all placeholders
            all_placeholders = [shape for shape in slide.shapes if hasattr(shape, 'placeholder_format')]
            
            # Look for subtitle placeholder specifically
            subtitle = None
            for shape in all_placeholders:
                # In PowerPoint, subtitle can be placeholder type 2 or 4 (SUBTITLE)
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type in (2, 4):
                    subtitle = shape
                    break
            
            # If subtitle placeholder found
            if subtitle:
                # Clear any existing content including the "Click to add subtitle" text
                subtitle.text = ""
                
                # Build subtitle text
                if headline.subtitle:
                    subtitle_text = headline.subtitle + "\n" + headline.date.strftime('%B %d, %Y')
                else:
                    subtitle_text = headline.date.strftime('%B %d, %Y')
                
                # Set the text directly
                subtitle.text = subtitle_text
                
                # Format the text
                text_frame = subtitle.text_frame
                configure_text_frame(text_frame)
                
                # Apply formatting to individual paragraphs if needed
                for i, paragraph in enumerate(text_frame.paragraphs):
                    if i == 0 and headline.subtitle:
                        paragraph.font.size = Pt(20)
                    else:
                        paragraph.font.size = Pt(18)
                    paragraph.alignment = PP_ALIGN.CENTER
                    if i < len(text_frame.paragraphs) - 1:
                        paragraph.space_after = Pt(12)
            else:
                raise AttributeError("No subtitle placeholder found")
        except (AttributeError, IndexError) as e:
            # Create separate textboxes for subtitle and date
            left = Inches(1)
            top = Inches(2.5)
            width = Inches(8)
            
            if headline.subtitle:
                # Subtitle in its own box
                subtitle_box = slide.shapes.add_textbox(left, top, width, height=Inches(1))
                subtitle_frame = subtitle_box.text_frame
                configure_text_frame(subtitle_frame)
                subtitle_frame.text = headline.subtitle
                subtitle_frame.paragraphs[0].font.size = Pt(20)
                subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Date in its own box, positioned below subtitle
            date_top = top + Inches(1.2)  # Position below subtitle
            date_box = slide.shapes.add_textbox(left, date_top, width, height=Inches(0.8))
            date_frame = date_box.text_frame
            configure_text_frame(date_frame)
            date_frame.text = headline.date.strftime('%B %d, %Y')
            date_frame.paragraphs[0].font.size = Pt(18)
            date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    except Exception as e:
        # Fall back to blank slide if all else fails
        logger.warning(f"Warning: {e}. Using blank slide for title.")
        slide = prs.slides.add_slide(prs.slide_layouts[-1])  # Usually blank layout
        
        # Add title textbox
        left = Inches(1)
        top = Inches(1)
        width = Inches(8)
        height = Inches(1.2)
        
        headline = Headline(**content.headline)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = title_box.text_frame
        configure_text_frame(text_frame)
        text_frame.text = headline.title
        text_frame.paragraphs[0].font.size = Pt(44)
        text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 112, 192)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Create separate textboxes for subtitle and date
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        
        if headline.subtitle:
            # Subtitle in its own box
            subtitle_box = slide.shapes.add_textbox(left, top, width, height=Inches(1))
            subtitle_frame = subtitle_box.text_frame
            configure_text_frame(subtitle_frame)
            subtitle_frame.text = headline.subtitle
            subtitle_frame.paragraphs[0].font.size = Pt(20)
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Date in its own box, positioned below subtitle
        date_top = top + Inches(1.2)  # Position below subtitle
        date_box = slide.shapes.add_textbox(left, date_top, width, height=Inches(0.8))
        date_frame = date_box.text_frame
        configure_text_frame(date_frame)
        date_frame.text = headline.date.strftime('%B %d, %Y')
        date_frame.paragraphs[0].font.size = Pt(18)
        date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add company name if provided
    if content.company:
        left = Inches(0.5)
        top = Inches(6)
        width = Inches(9)
        height = Inches(0.5)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        configure_text_frame(text_frame)
        text_frame.text = content.company
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.paragraphs[0].font.size = Pt(12)
        text_frame.paragraphs[0].font.italic = True
    
    return slide


def process_bullet_point(p, item):
    """Process a bullet point item which can be a string or BulletPoint object"""
    if isinstance(item, str):
        p.text = item
        p.level = 0
    else:  # BulletPoint object or dict
        bullet_point = item if isinstance(item, BulletPoint) else BulletPoint(**item)
        p.text = bullet_point.text
        p.level = bullet_point.level
        p.font.bold = bullet_point.bold
        p.font.italic = bullet_point.italic
        p.font.underline = bullet_point.underline
    
    # Set a reasonable font size for bullet points based on level
    if p.level == 0:
        p.font.size = Pt(24)
    elif p.level == 1:
        p.font.size = Pt(20)
    else:
        p.font.size = Pt(18)


def create_content_slide(prs, slide_content_dict):
    """Create a content slide with bullets, image, chart, or table"""
    # Fix for "bullets" vs "bullet_points" inconsistency
    if "bullets" in slide_content_dict and "bullet_points" not in slide_content_dict:
        slide_content_dict["bullet_points"] = slide_content_dict.pop("bullets")

    # Convert dict to ContentSlide object for validation
    slide_content = ContentSlide(**slide_content_dict)
    
    # Choose an appropriate layout or fall back to blank slide if needed
    try:
        # Find a suitable layout based on content
        if slide_content.chart_data:
            slide_layout = prs.slide_layouts[5]  # Title and Content layout
        elif slide_content.table_data:
            slide_layout = prs.slide_layouts[5]  # Title and Content layout
        elif slide_content.image_path:
            slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[1]
        else:
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
    except (IndexError, AttributeError):
        # Fall back to the first available layout if specified ones aren't available
        slide_layout = prs.slide_layouts[0] if prs.slide_layouts else None

    # Create slide
    if slide_layout:
        slide = prs.slides.add_slide(slide_layout)
    else:
        # If no layouts are available, report an error
        logger.warning("Warning: No slide layouts available, creating blank slide")
        slide = prs.slides.add_slide(prs.slide_layouts[-1])  # Usually blank layout
    
    # Set title
    try:
        title = slide.shapes.title
        title.text = slide_content.title
        configure_text_frame(title.text_frame)
    except (AttributeError, IndexError):
        # Create a textbox for title if no title placeholder
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(0.8)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = title_box.text_frame
        configure_text_frame(text_frame)
        text_frame.text = slide_content.title
        text_frame.paragraphs[0].font.size = Pt(32)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    content_left = Inches(0.5)
    content_top = Inches(1.5)
    content_width = Inches(9)
    content_height = Inches(5)
    
    # Adjust layout if we have chart or image
    if slide_content.chart_data or slide_content.image_path:
        if slide_content.bullet_points:
            # If we have both bullet points and chart/image, adjust layout
            content_width = Inches(4.5)  # Make room for image/chart on right
    
    # Handle simple content field if provided (new functionality)
    content_field = getattr(slide_content, 'content', None)
    if content_field:
        try:
            # Try to find content placeholder
            content_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 7:  # Body placeholder
                    content_placeholder = shape
                    break
            
            if content_placeholder:
                text_frame = content_placeholder.text_frame
                text_frame.clear()  # Clear existing text
                configure_text_frame(text_frame)
                text_frame.text = content_field
            else:
                raise AttributeError("No content placeholder found")
        except (IndexError, AttributeError):
            # If no placeholder available, add a textbox
            textbox = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
            text_frame = textbox.text_frame
            configure_text_frame(text_frame)
            text_frame.text = content_field
            text_frame.paragraphs[0].font.size = Pt(18)
    
    # Add bullet points if provided
    if slide_content.bullet_points:
        try:
            # Try to find content placeholder
            content_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 7:  # Content placeholder
                    content_placeholder = shape
                    break
            
            if content_placeholder:
                text_frame = content_placeholder.text_frame
                text_frame.clear()  # Clear existing text
                configure_text_frame(text_frame)
                
                for item in slide_content.bullet_points:
                    p = text_frame.add_paragraph()
                    process_bullet_point(p, item)
            else:
                raise AttributeError("No content placeholder found")
        except (IndexError, AttributeError):
            # If no placeholder available, add a textbox
            textbox = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
            text_frame = textbox.text_frame
            configure_text_frame(text_frame)
            
            for item in slide_content.bullet_points:
                p = text_frame.add_paragraph()
                process_bullet_point(p, item)
    
    # Add image if provided
    if slide_content.image_path and os.path.exists(slide_content.image_path):
        image_left = Inches(5.5) if slide_content.bullet_points else Inches(2.5)
        image_top = Inches(2)
        image_width = Inches(4) if slide_content.bullet_points else Inches(6)
        
        try:
            # Try to find picture placeholder
            picture_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 18:  # Picture placeholder
                    picture_placeholder = shape
                    break
            
            if picture_placeholder:
                picture_placeholder.insert_picture(slide_content.image_path)
            else:
                raise AttributeError("No picture placeholder found")
        except (IndexError, AttributeError):
            # If no placeholder, add picture directly
            slide.shapes.add_picture(
                slide_content.image_path, 
                image_left, 
                image_top, 
                width=image_width
            )
    
    # Add chart if provided
    if slide_content.chart_data:
        chart_data_obj = ChartData(**slide_content.chart_data)
        chart_data = CategoryChartData()
        chart_data.categories = chart_data_obj.categories
        
        for series_name, values in chart_data_obj.series.items():
            chart_data.add_series(series_name, values)
        
        # Determine chart type
        chart_type_dict = {
            "BAR": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "LINE": XL_CHART_TYPE.LINE,
            "PIE": XL_CHART_TYPE.PIE,
            "AREA": XL_CHART_TYPE.AREA,
        }
        chart_type = chart_type_dict.get(
            chart_data_obj.chart_type.upper(), 
            XL_CHART_TYPE.COLUMN_CLUSTERED
        )
        
        # Set chart dimensions based on whether we have bullet points
        chart_left = Inches(5.5) if slide_content.bullet_points else Inches(1.5)
        chart_top = Inches(2)
        chart_width = Inches(4) if slide_content.bullet_points else Inches(7)
        chart_height = Inches(4)
        
        try:
            # Try to find chart placeholder
            chart_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 7:  # Content placeholder
                    chart_placeholder = shape
                    break
            
            if chart_placeholder:
                chart = chart_placeholder.insert_chart(chart_type, chart_data)
                chart.chart_title.text_frame.text = chart_data_obj.title
                configure_text_frame(chart.chart_title.text_frame)
            else:
                raise AttributeError("No chart placeholder found")
        except (IndexError, AttributeError):
            # If no placeholder, add chart directly
            chart = slide.shapes.add_chart(
                chart_type, chart_left, chart_top, chart_width, chart_height, chart_data
            ).chart
            chart.chart_title.text_frame.text = chart_data_obj.title
            configure_text_frame(chart.chart_title.text_frame)
    
    # Add table if provided
    if slide_content.table_data:
        table_data_obj = TableData(**slide_content.table_data)
        rows = len(table_data_obj.rows) + 1  # +1 for header
        cols = len(table_data_obj.headers)
        
        # Calculate appropriate table size and position
        table_left = Inches(0.5)
        table_top = content_top
        table_width = Inches(9)
        row_height = min(0.4, 3.0 / rows)  # Limit row height for many rows
        table_height = Inches(row_height * rows)
        
        # Add table shape
        table = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height).table
        
        # Set consistent column widths - Fixed to avoid the '_Column' object has no attribute 'cells' error
        col_width = table_width.inches / cols
        for i in range(cols):
            for row_idx in range(rows):
                table.cell(row_idx, i).width = Inches(col_width)
        
        # Add headers
        for i, header in enumerate(table_data_obj.headers):
            cell = table.cell(0, i)
            cell.text = header
            # Format header text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.bold = True
                paragraph.font.size = Pt(14)
                paragraph.alignment = PP_ALIGN.CENTER
            configure_text_frame(cell.text_frame, auto_size=False)
        
        # Add rows
        for row_idx, row_data in enumerate(table_data_obj.rows):
            for col_idx, cell_text in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)  # +1 to skip header row
                cell.text = cell_text
                # Format cell text
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(12)
                    paragraph.alignment = PP_ALIGN.CENTER
                configure_text_frame(cell.text_frame, auto_size=False)
    
    # Add notes if provided
    if slide_content.notes:
        try:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_content.notes
            configure_text_frame(text_frame)
        except Exception as e:
            logger.warning(f"Warning: Could not add notes to slide: {e}")
    
    return slide


def add_footer(slide, text):
    """Add footer to slide"""
    left = Inches(0.5)
    top = Inches(6.9)  # Moved up slightly to stay within slide
    width = Inches(9)
    height = Inches(0.4)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = text
    text_frame.paragraphs[0].font.size = Pt(10)
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)  # Grey color


def set_presentation_properties(prs, content):
    """Set presentation metadata properties"""
    # Basic properties
    prs.core_properties.title = content.title
    prs.core_properties.author = content.author
    prs.core_properties.comments = f"Presentation created with python-pptx"
    prs.core_properties.category = "Presentation"
    
    # Keywords
    if content.keywords:
        prs.core_properties.keywords = ", ".join(content.keywords)
    
    # Other properties
    prs.core_properties.revision = content.revision
    prs.core_properties.last_modified_by = content.author
    
    # Set subject to headline title if headline exists
    if hasattr(content, 'headline') and content.headline is not None:
        headline = Headline(**content.headline)
        prs.core_properties.subject = headline.title
    else:
        # If no headline, use the presentation title as subject
        prs.core_properties.subject = content.title
        
    prs.core_properties.content_status = "Final"


def load_content_from_json(json_file_path):
    """Load presentation content from a JSON file"""
    try:
        with open(json_file_path, 'r', encoding='utf-8') as f:
            content_data = json.load(f)
            
        # Parse date string in headline if present
        if 'headline' in content_data and 'date' in content_data['headline']:
            date_str = content_data['headline']['date']
            if isinstance(date_str, str):
                content_data['headline']['date'] = Headline.parse_date(date_str)
                
        # Validate using Pydantic model
        content = PresentationContent(**content_data)
        return content
    except (FileNotFoundError, json.JSONDecodeError, ValidationError) as e:
        logger.error(f"Error loading content from JSON: {e}")
        raise


def load_content_from_dict(content_dict):
    """Load presentation content from a dictionary"""
    try:
        # Ensure headline exists
        if 'headline' not in content_dict or content_dict['headline'] is None:
            # Create default headline using the title
            title = content_dict.get('title', 'Presentation')
            content_dict['headline'] = {
                'title': title,
                'date': datetime.now()
            }
        # Parse date string in headline if present
        elif 'date' in content_dict['headline']:
            date_str = content_dict['headline']['date']
            if isinstance(date_str, str):
                content_dict['headline']['date'] = Headline.parse_date(date_str)
                
        # Process slides array to handle simpler formats
        if 'slides' in content_dict and isinstance(content_dict['slides'], list):
            # Make sure each slide conforms to expected format
            for i, slide in enumerate(content_dict['slides']):
                # No changes needed, our ContentSlide model now handles 'content' field
                pass
                
        # Validate using Pydantic model
        content = PresentationContent(**content_dict)
        return content
    except ValidationError as e:
        logger.error(f"Error validating content data: {e}")
        raise


def create_presentation(content, output_path="presentation.pptx"):
    """Create full presentation from content model"""
    # Create a new presentation
    prs = Presentation()
    
    # Set standard slide size (16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)  # Increased from 5.625 for better content fit
    
    # Set presentation properties
    set_presentation_properties(prs, content)
    
    # Create title slide
    title_slide = create_title_slide(prs, content)
    
    # Create content slides
    for slide_content in content.slides:
        slide = create_content_slide(prs, slide_content)
        
        # Add footer if required
        if content.include_footer:
            add_footer(slide, f"{content.title} | {datetime.now().strftime('%B %Y')}")
    
    # Save the presentation
    prs.save(output_path)
    logger.info(f"Presentation saved to {output_path}")
    
    return prs


# PPTXGenerator class for CrewAI Tool Interface
class PPTXGenerator:
    """Class to generate PowerPoint presentations."""

    def __init__(self, output_dir: str = "./", template_path: Optional[str] = None):
        """Initialize the generator.

        Args:
            output_dir: Directory to save the presentation
            template_path: Path to a PPTX template file (optional)
        """
        self.output_dir = output_dir
        self.template_path = template_path

    def generate(self, content: str, title: Optional[str] = None) -> Dict[str, str]:
        """Generate a presentation from content.

        Args:
            content: The raw content for the slides (text or JSON)
            title: Title of the presentation

        Returns:
            Dict with file paths to the presentation
        """
        # Try to parse content as JSON, otherwise treat as plain text
        try:
            # If content can be parsed as JSON, use structured approach
            structured_content = json.loads(content) if isinstance(content, str) else content
            return self.generate_from_json(structured_content, title)
        except (json.JSONDecodeError, ValueError):
            # Handle as simple text content
            prs = self._create_presentation()
            
            # Create title slide if we have a title
            if title:
                self._add_title_slide(prs, title)
            
            # Create a content slide
            self._add_content_slide(prs, "Content", content)
            
            return self._save_presentation(prs, title or "Presentation")
        except Exception as e:
            logging.error(f"Error in generate: {e}")
            logging.error(traceback.format_exc())
            raise

    def generate_from_json(self, content: Dict[str, Any], title: Optional[str] = None) -> Dict[str, str]:
        """Generate a presentation from structured JSON content.

        Args:
            content: Structured content dictionary
            title: Title of the presentation

        Returns:
            Dict with file paths to the presentation
        """
        try:
            # Extract basic metadata
            presentation_title = title or content.get("title", "Presentation")
            author = content.get("author", "")
            description = content.get("description", "")
            company = content.get("company", "")
            keywords = content.get("keywords", [])
            
            # Ensure headline exists
            if "headline" not in content or content["headline"] is None:
                # Create default headline
                content["headline"] = {
                    "title": presentation_title,
                    "date": datetime.now()
                }
            
            # Create presentation
            prs = self._create_presentation()
            
            # Set presentation metadata
            if hasattr(prs, 'core_properties'):
                prs.core_properties.title = presentation_title
                prs.core_properties.author = author
                prs.core_properties.comments = description
                if keywords:
                    prs.core_properties.keywords = ', '.join(keywords) if isinstance(keywords, list) else keywords
            
            # Create title slide
            if "headline" in content:
                # Create title slide with headline data
                headline = content["headline"]
                title_slide_layout = prs.slide_layouts[0]
                title_slide = prs.slides.add_slide(title_slide_layout)
                
                if hasattr(title_slide, "shapes") and hasattr(title_slide.shapes, "title"):
                    title_shape = title_slide.shapes.title
                    if title_shape:
                        title_shape.text = headline.get("title", presentation_title)
                
                # Add subtitle if present
                for shape in title_slide.placeholders:
                    if shape.placeholder_format.type == 2:  # Subtitle placeholder
                        if "subtitle" in headline:
                            shape.text = headline["subtitle"]
                        # Add date if available
                        if "date" in headline:
                            try:
                                date_str = headline["date"]
                                if isinstance(date_str, str):
                                    # Try to parse date
                                    date_format = '%B %d, %Y'
                                    try:
                                        # Try ISO format first
                                        date_obj = datetime.fromisoformat(date_str.replace('Z', '+00:00'))
                                        p = shape.text_frame.add_paragraph()
                                        p.text = date_obj.strftime(date_format)
                                    except ValueError:
                                        # Use as-is
                                        p = shape.text_frame.add_paragraph()
                                        p.text = date_str
                            except Exception as e:
                                logging.warning(f"Error parsing date: {e}")
            else:
                # Simple title slide
                self._add_title_slide(prs, presentation_title, description)
            
            # Create content slides
            if "slides" in content and isinstance(content["slides"], list):
                slides_data = content["slides"]
                for slide_data in slides_data:
                    if isinstance(slide_data, dict):
                        slide_title = slide_data.get("title", "")
                        
                        # Handle both "bullets" and "bullet_points" for compatibility
                        if "bullets" in slide_data and "bullet_points" not in slide_data:
                            slide_data["bullet_points"] = slide_data.pop("bullets")
                        
                        # Process slide based on content type
                        if "bullet_points" in slide_data:
                            # Bullet points slide
                            self._add_bullet_slide(prs, slide_title, slide_data["bullet_points"], slide_data.get("notes"))
                        elif "chart_data" in slide_data:
                            # Chart slide
                            self._add_chart_slide(prs, slide_title, slide_data["chart_data"], slide_data.get("notes"))
                        elif "table_data" in slide_data:
                            # Table slide
                            self._add_table_slide(prs, slide_title, slide_data["table_data"], slide_data.get("notes"))
                        elif "content" in slide_data:
                            # Simple content slide
                            self._add_content_slide(prs, slide_title, slide_data["content"], slide_data.get("notes"))
                        else:
                            # Empty slide with just a title
                            self._add_content_slide(prs, slide_title, "", slide_data.get("notes"))
            else:
                logging.warning("No 'slides' list found in content")
            
            # Save the presentation
            return self._save_presentation(prs, presentation_title)
        except Exception as e:
            logging.error(f"Error in generate_from_json: {e}")
            logging.error(traceback.format_exc())
            raise

    def _create_presentation(self) -> Presentation:
        """Create a new presentation object.
        
        Returns:
            A new Presentation object
        """
        if self.template_path and os.path.exists(self.template_path):
            return Presentation(self.template_path)
        return Presentation()

    def _add_title_slide(self, prs: Presentation, title: str, subtitle: str = "") -> None:
        """Add a title slide to the presentation.
        
        Args:
            prs: Presentation object
            title: Title text
            subtitle: Optional subtitle text
        """
        try:
            slide_layout = prs.slide_layouts[0]  # Title slide layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            if hasattr(slide, "shapes") and hasattr(slide.shapes, "title") and slide.shapes.title:
                slide.shapes.title.text = title
            
            # Set subtitle if there's a placeholder for it
            subtitle_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 2:  # Subtitle placeholder
                    subtitle_placeholder = shape
                    break
                    
            if subtitle_placeholder:
                subtitle_placeholder.text = subtitle or datetime.now().strftime("%B %d, %Y")
        except Exception as e:
            logging.warning(f"Error adding title slide: {e}")

    def _add_content_slide(self, prs: Presentation, title: str, content: str, notes: Optional[str] = None) -> None:
        """Add a simple content slide with title and text content.
        
        Args:
            prs: Presentation object
            title: Slide title
            content: Text content
            notes: Optional speaker notes
        """
        try:
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            if hasattr(slide, "shapes") and hasattr(slide.shapes, "title") and slide.shapes.title:
                slide.shapes.title.text = title
            
            # Set content if we have a content placeholder
            body_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type in (2, 7):  # Body placeholder
                    body_placeholder = shape
                    break
                    
            if body_placeholder and content:
                body_placeholder.text = content
            
            # Add notes if provided
            if notes and hasattr(slide, "notes_slide"):
                slide.notes_slide.notes_text_frame.text = notes
                
            return slide
        except Exception as e:
            logging.warning(f"Error adding content slide: {e}")
            return None

    def _add_bullet_slide(self, prs: Presentation, title: str, bullets: List[Union[str, Dict]], notes: Optional[str] = None) -> None:
        """Add a slide with bullet points.
        
        Args:
            prs: Presentation object
            title: Slide title
            bullets: List of bullet points (strings or dicts with formatting)
            notes: Optional speaker notes
        """
        try:
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            if hasattr(slide, "shapes") and hasattr(slide.shapes, "title") and slide.shapes.title:
                slide.shapes.title.text = title
            
            # Find content placeholder
            content_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 7:  # Content placeholder
                    content_placeholder = shape
                    break
                    
            if not content_placeholder:
                # If no content placeholder, add a textbox
                left = Inches(0.5)
                top = Inches(1.5)
                width = Inches(9)
                height = Inches(5)
                
                content_placeholder = slide.shapes.add_textbox(left, top, width, height)
                
            # Add bullet points
            tf = content_placeholder.text_frame
            tf.clear()  # Clear any existing text
            
            for bullet in bullets:
                p = tf.add_paragraph()
                
                # Handle different bullet point formats
                if isinstance(bullet, str):
                    p.text = bullet
                    p.level = 0
                elif isinstance(bullet, dict):
                    p.text = bullet.get("text", "")
                    
                    # Apply bullet level if specified
                    if "level" in bullet:
                        p.level = min(bullet["level"], 4)  # Max level 4
                    
                    # Apply formatting if specified
                    if bullet.get("bold", False):
                        p.font.bold = True
                    if bullet.get("italic", False):
                        p.font.italic = True
            
            # Add notes if provided
            if notes and hasattr(slide, "notes_slide"):
                slide.notes_slide.notes_text_frame.text = notes
                
            return slide
        except Exception as e:
            logging.warning(f"Error adding bullet slide: {e}")
            return None

    def _add_chart_slide(self, prs: Presentation, title: str, chart_data: Dict[str, Any], notes: Optional[str] = None) -> None:
        """Add a slide with a chart.
        
        Args:
            prs: Presentation object
            title: Slide title
            chart_data: Chart configuration
            notes: Optional speaker notes
        """
        try:
            slide_layout = prs.slide_layouts[5]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            if hasattr(slide, "shapes") and hasattr(slide.shapes, "title") and slide.shapes.title:
                slide.shapes.title.text = title
            
            # Prepare chart data
            chart_title = chart_data.get("title", "Chart")
            
            # Handle different chart data structures
            chart_type_str = chart_data.get("chart_type", chart_data.get("type", "BAR"))
            chart_type_str = chart_type_str.upper() if isinstance(chart_type_str, str) else "BAR"
            
            # Convert chart type to python-pptx enum
            chart_type_dict = {
                "BAR": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "COLUMN": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "LINE": XL_CHART_TYPE.LINE,
                "PIE": XL_CHART_TYPE.PIE,
                "AREA": XL_CHART_TYPE.AREA,
                "SCATTER": XL_CHART_TYPE.XY_SCATTER,
                "RADAR": XL_CHART_TYPE.RADAR,
                "DOUGHNUT": XL_CHART_TYPE.DOUGHNUT,
            }
            xl_chart_type = chart_type_dict.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)
            
            # Get categories and datasets
            categories = []
            
            # Extract categories
            if "categories" in chart_data:
                categories = chart_data["categories"]
            elif "data" in chart_data and "labels" in chart_data["data"]:
                categories = chart_data["data"]["labels"]
            
            # Create chart data object
            cd = CategoryChartData()
            cd.categories = categories
            
            # Get series data
            series_data = {}
            if "series" in chart_data:
                # Original format
                series_data = chart_data["series"]
            elif "data" in chart_data and "datasets" in chart_data["data"]:
                # Format seen in examples
                for dataset in chart_data["data"]["datasets"]:
                    if "label" in dataset and "data" in dataset:
                        series_data[dataset["label"]] = dataset["data"]
            
            # Add series data
            for series_name, values in series_data.items():
                cd.add_series(series_name, values)
            
            # Add chart to slide (position and size)
            chart_left = Inches(2)
            chart_top = Inches(2)
            chart_width = Inches(6)
            chart_height = Inches(4)
            
            chart = slide.shapes.add_chart(
                xl_chart_type, chart_left, chart_top, chart_width, chart_height, cd
            ).chart
            
            chart.has_title = True
            chart.chart_title.text_frame.text = chart_title
            
            # Add notes if provided
            if notes and hasattr(slide, "notes_slide"):
                slide.notes_slide.notes_text_frame.text = notes
                
            return slide
        except Exception as e:
            logging.warning(f"Error adding chart slide: {e}")
            return None

    def _add_table_slide(self, prs: Presentation, title: str, table_data: Dict[str, Any], notes: Optional[str] = None) -> None:
        """Add a slide with a table.
        
        Args:
            prs: Presentation object
            title: Slide title
            table_data: Table configuration with headers and rows
            notes: Optional speaker notes
        """
        try:
            # Add slide with title
            slide_layout = prs.slide_layouts[5]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            if hasattr(slide, "shapes") and hasattr(slide.shapes, "title") and slide.shapes.title:
                slide.shapes.title.text = title
                configure_text_frame(slide.shapes.title.text_frame)
            
            # Extract headers and rows
            headers = table_data.get("headers", table_data.get("columns", []))
            rows = table_data.get("rows", [])
            
            # Calculate table dimensions
            if not headers or not rows:
                logging.warning(f"Missing table data for slide: {title}")
                return None
                
            num_cols = len(headers)
            num_rows = len(rows) + 1  # +1 for header row
            
            # Calculate position and size
            table_left = Inches(0.5)
            table_top = Inches(1.8)
            table_width = Inches(9)
            
            # Adjust row height based on content
            row_height = 0.4  # Default height in inches
            if "row_height" in table_data:
                row_height = table_data["row_height"]
            else:
                # Calculate reasonable height based on number of rows
                available_height = 5.0  # Available height in inches
                row_height = min(0.6, available_height / num_rows)
            
            table_height = Inches(row_height * num_rows)
            
            # Apply custom column widths if specified
            column_widths = table_data.get("column_widths")
            
            # Create table
            table = slide.shapes.add_table(
                num_rows, num_cols, table_left, table_top, table_width, table_height
            ).table
            
            # Set column widths if specified
            if column_widths and len(column_widths) == num_cols:
                total_width = sum(column_widths)
                table_width_inches = table_width.inches
                
                for i, width in enumerate(column_widths):
                    # Convert relative width to inches
                    col_width = (width / total_width) * table_width_inches
                    table.columns[i].width = Inches(col_width)
            
            # Add headers with formatting
            for i, header in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = str(header)
                
                # Format header text
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0, 112, 192)  # Blue
                
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White
                    paragraph.alignment = PP_ALIGN.CENTER
                    paragraph.font.size = Pt(14)
            
            # Add rows with formatting
            for i, row_data in enumerate(rows):
                for j, cell_text in enumerate(row_data):
                    if j < num_cols:  # Stay within bounds
                        cell = table.cell(i + 1, j)  # +1 to skip header row
                        cell.text = str(cell_text)
                        
                        # Format cell text
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.size = Pt(12)
                            paragraph.alignment = PP_ALIGN.CENTER
                        
                        # Apply alternating row colors
                        if i % 2 == 1:  # Every other row
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(242, 242, 242)  # Light gray
            
            # Apply cell-specific formatting if provided
            if "cell_formatting" in table_data:
                for cell_format in table_data["cell_formatting"]:
                    row = cell_format.get("row", 0)
                    col = cell_format.get("col", 0)
                    
                    if 0 <= row < num_rows and 0 <= col < num_cols:
                        cell = table.cell(row, col)
                        
                        # Apply text formatting
                        if "text" in cell_format:
                            cell.text = str(cell_format["text"])
                        
                        # Apply cell formatting
                        if "bold" in cell_format and cell_format["bold"]:
                            for paragraph in cell.text_frame.paragraphs:
                                paragraph.font.bold = True
                        
                        if "font_size" in cell_format:
                            for paragraph in cell.text_frame.paragraphs:
                                paragraph.font.size = Pt(cell_format["font_size"])
                        
                        if "alignment" in cell_format:
                            align_map = {
                                "left": PP_ALIGN.LEFT,
                                "center": PP_ALIGN.CENTER,
                                "right": PP_ALIGN.RIGHT,
                                "justify": PP_ALIGN.JUSTIFY
                            }
                            if cell_format["alignment"].lower() in align_map:
                                for paragraph in cell.text_frame.paragraphs:
                                    paragraph.alignment = align_map[cell_format["alignment"].lower()]
                        
                        # Apply cell fill
                        if "fill_color" in cell_format:
                            cell.fill.solid()
                            set_shape_fill_color(cell, cell_format["fill_color"])
            
            # Add notes if provided
            if notes and hasattr(slide, "notes_slide"):
                slide.notes_slide.notes_text_frame.text = notes
                
            return slide
        except Exception as e:
            logging.warning(f"Error adding table slide: {e}")
            return None

    def _add_shapes_slide(self, prs: Presentation, title: str, shapes_data: List[Dict], notes: Optional[str] = None) -> None:
        """Add a slide with custom shapes.
        
        Args:
            prs: Presentation object
            title: Slide title
            shapes_data: List of shape configurations
            notes: Optional speaker notes
        """
        try:
            slide_layout = prs.slide_layouts[5]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            if hasattr(slide, "shapes") and hasattr(slide.shapes, "title") and slide.shapes.title:
                slide.shapes.title.text = title
                configure_text_frame(slide.shapes.title.text_frame)
            
            # Add each shape to the slide
            for shape_data in shapes_data:
                shape_type = shape_data.get("type", "rectangle").lower()
                
                # Extract common properties
                x = shape_data.get("x", 1.0)
                y = shape_data.get("y", 1.0)
                width = shape_data.get("width", 2.0)
                height = shape_data.get("height", 1.0)
                
                # Create the appropriate shape
                if shape_type == "rectangle":
                    shape = slide.shapes.add_shape(
                        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                        Inches(x), Inches(y), Inches(width), Inches(height)
                    )
                elif shape_type == "oval" or shape_type == "ellipse":
                    shape = slide.shapes.add_shape(
                        MSO_AUTO_SHAPE_TYPE.OVAL,
                        Inches(x), Inches(y), Inches(width), Inches(height)
                    )
                elif shape_type == "rounded_rectangle":
                    shape = slide.shapes.add_shape(
                        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                        Inches(x), Inches(y), Inches(width), Inches(height)
                    )
                elif shape_type == "triangle":
                    shape = slide.shapes.add_shape(
                        MSO_AUTO_SHAPE_TYPE.TRIANGLE,
                        Inches(x), Inches(y), Inches(width), Inches(height)
                    )
                elif shape_type == "pentagon":
                    shape = slide.shapes.add_shape(
                        MSO_AUTO_SHAPE_TYPE.PENTAGON,
                        Inches(x), Inches(y), Inches(width), Inches(height)
                    )
                elif shape_type == "hexagon":
                    shape = slide.shapes.add_shape(
                        MSO_AUTO_SHAPE_TYPE.HEXAGON,
                        Inches(x), Inches(y), Inches(width), Inches(height)
                    )
                elif shape_type == "arrow":
                    shape = slide.shapes.add_shape(
                        MSO_AUTO_SHAPE_TYPE.ARROW,
                        Inches(x), Inches(y), Inches(width), Inches(height)
                    )
                elif shape_type == "star":
                    shape = slide.shapes.add_shape(
                        MSO_AUTO_SHAPE_TYPE.STAR_5_POINTS,
                        Inches(x), Inches(y), Inches(width), Inches(height)
                    )
                elif shape_type == "freeform" and "vertices" in shape_data:
                    # Create freeform shape
                    shape = create_freeform_shape(
                        slide, 
                        shape_data["vertices"],
                        fill_color=shape_data.get("fill_color"),
                        line_color=shape_data.get("line_color"),
                        line_width=shape_data.get("line_width")
                    )
                else:
                    # Default to rectangle
                    shape = slide.shapes.add_shape(
                        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                        Inches(x), Inches(y), Inches(width), Inches(height)
                    )
                
                # Apply fill formatting
                fill_type = shape_data.get("fill_type", "solid")
                if fill_type == "solid":
                    if "fill_color" in shape_data:
                        shape.fill.solid()
                        set_shape_fill_color(shape, shape_data["fill_color"])
                elif fill_type == "gradient" and "gradient_start" in shape_data and "gradient_end" in shape_data:
                    apply_gradient_fill(
                        shape,
                        shape_data["gradient_start"],
                        shape_data["gradient_end"],
                        gradient_type=shape_data.get("gradient_type", "linear"),
                        angle=shape_data.get("gradient_angle", 90)
                    )
                elif fill_type == "none":
                    shape.fill.background()
                
                # Apply line formatting
                if "line_color" in shape_data:
                    set_shape_line_color(shape, shape_data["line_color"])
                
                if "line_width" in shape_data:
                    shape.line.width = Pt(shape_data["line_width"])
                
                if "line_dash" in shape_data:
                    set_line_dash_style(shape.line, shape_data["line_dash"])
                
                # Apply text if provided
                if "text" in shape_data:
                    text_data = shape_data["text"]
                    
                    if isinstance(text_data, str):
                        shape.text = text_data
                        
                        # Apply basic text formatting
                        if "text_align" in shape_data:
                            align_map = {
                                "left": PP_ALIGN.LEFT,
                                "center": PP_ALIGN.CENTER,
                                "right": PP_ALIGN.RIGHT,
                                "justify": PP_ALIGN.JUSTIFY
                            }
                            if shape_data["text_align"].lower() in align_map:
                                for paragraph in shape.text_frame.paragraphs:
                                    paragraph.alignment = align_map[shape_data["text_align"].lower()]
                        
                        if "text_color" in shape_data:
                            for paragraph in shape.text_frame.paragraphs:
                                set_font_color(paragraph.font, shape_data["text_color"])
                        
                        if "font_bold" in shape_data:
                            for paragraph in shape.text_frame.paragraphs:
                                paragraph.font.bold = shape_data["font_bold"]
                        
                        if "font_size" in shape_data:
                            for paragraph in shape.text_frame.paragraphs:
                                paragraph.font.size = Pt(shape_data["font_size"])
                    elif isinstance(text_data, dict) or isinstance(text_data, TextFrame):
                        # Use text frame with rich formatting
                        shape.text = ""  # Clear default text
                        configure_text_frame(
                            shape.text_frame,
                            auto_size=text_data.get("auto_size", True),
                            word_wrap=text_data.get("word_wrap", True),
                            vertical_alignment=text_data.get("vertical_alignment"),
                            margins=text_data.get("margins")
                        )
                        
                        # Add paragraphs with rich formatting
                        for para_data in text_data.get("paragraphs", []):
                            add_text_with_mixed_formatting(shape.text_frame, para_data)
                
                # Apply hyperlink if provided
                if "hyperlink" in shape_data:
                    hyperlink_data = shape_data["hyperlink"]
                    url = hyperlink_data["url"] if isinstance(hyperlink_data, dict) else hyperlink_data
                    tooltip = hyperlink_data.get("tooltip") if isinstance(hyperlink_data, dict) else None
                    add_hyperlink(shape, url, tooltip)
                
                # Apply shadow if requested
                if shape_data.get("shadow", False):
                    shape.shadow.inherit = False
                    shape.shadow.visible = True
                    
                    # Configure shadow properties if provided
                    shadow_props = shape_data.get("shadow_properties", {})
                    if "blur_radius" in shadow_props:
                        shape.shadow.blur_radius = Pt(shadow_props["blur_radius"])
                    if "distance" in shadow_props:
                        shape.shadow.distance = Pt(shadow_props["distance"])
                    if "angle" in shadow_props:
                        shape.shadow.angle = shadow_props["angle"]
                
                # Apply rotation if specified
                if "rotation" in shape_data:
                    shape.rotation = shape_data["rotation"]
            
            # Add notes if provided
            if notes and hasattr(slide, "notes_slide"):
                slide.notes_slide.notes_text_frame.text = notes
                
            return slide
        except Exception as e:
            logging.warning(f"Error adding shapes slide: {e}")
            logging.error(traceback.format_exc())
            return None

    def _add_mixed_content_slide(self, prs: Presentation, title: str, slide_data: Dict[str, Any]) -> None:
        """Add a slide with multiple content types.
        
        Args:
            prs: Presentation object
            title: Slide title
            slide_data: Dictionary with multiple content types
        """
        try:
            slide_layout = prs.slide_layouts[5]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            if hasattr(slide, "shapes") and hasattr(slide.shapes, "title") and slide.shapes.title:
                slide.shapes.title.text = title
                configure_text_frame(slide.shapes.title.text_frame)
            
            # Track if we've created a content shape
            has_content_shape = False
            content_right = 9.0  # right edge in inches
            
            # Determine if we should split the slide or stack content
            has_chart = "chart_data" in slide_data
            has_table = "table_data" in slide_data
            has_bullet_points = "bullet_points" in slide_data or "bullets" in slide_data
            has_image = "image_path" in slide_data and os.path.exists(slide_data["image_path"])
            has_text_content = "content" in slide_data or "paragraphs" in slide_data
            
            # Config to split horizontally (bullet points on left, chart/image on right)
            split_horizontally = (
                (has_bullet_points or has_text_content) and (has_chart or has_image) and 
                not (has_bullet_points and has_chart and has_image and has_table)
            )
            
            # Calculate layout
            left_column_width = 4.5 if split_horizontally else 9.0
            right_column_left = 5.0 if split_horizontally else 0.0
            right_column_width = 4.0 if split_horizontally else 9.0
            
            # Create content areas
            if has_bullet_points:
                # Add bullet points (left side if split, full width otherwise)
                left = Inches(0.5)
                top = Inches(1.8)
                width = Inches(left_column_width)
                height = Inches(5)
                
                # Get bullet points
                bullet_points = slide_data.get("bullet_points", slide_data.get("bullets", []))
                
                # Create textbox for bullet points
                textbox = slide.shapes.add_textbox(left, top, width, height)
                tf = textbox.text_frame
                configure_text_frame(tf)
                
                # Add bullet points
                for bullet in bullet_points:
                    p = tf.add_paragraph()
                    
                    # Handle different bullet point formats
                    if isinstance(bullet, str):
                        p.text = bullet
                        p.level = 0
                        # Set default font size
                        p.font.size = Pt(18 if split_horizontally else 20)
                    elif isinstance(bullet, dict):
                        if "text" in bullet:
                            p.text = bullet["text"]
                        else:
                            p.text = str(bullet)
                        
                        # Apply bullet level if specified
                        if "level" in bullet:
                            p.level = min(int(bullet["level"]), 4)  # Max level 4
                        
                        # Apply formatting if specified
                        if bullet.get("bold", False):
                            p.font.bold = True
                        if bullet.get("italic", False):
                            p.font.italic = True
                        
                        # Apply font size based on level or explicit setting
                        if "font_size" in bullet:
                            p.font.size = Pt(bullet["font_size"])
                        else:
                            base_size = 16 if split_horizontally else 18
                            p.font.size = Pt(base_size - (p.level * 2))
                
                has_content_shape = True
            
            # Add simple text content if provided
            if "content" in slide_data and not has_bullet_points:
                # Use the same layout as bullet points would use
                left = Inches(0.5)
                top = Inches(1.8)
                width = Inches(left_column_width)
                height = Inches(5)
                
                # Create textbox
                textbox = slide.shapes.add_textbox(left, top, width, height)
                tf = textbox.text_frame
                configure_text_frame(tf)
                tf.text = slide_data["content"]
                tf.paragraphs[0].font.size = Pt(18 if split_horizontally else 20)
                
                has_content_shape = True
            
            # Add rich text paragraphs if provided
            elif "paragraphs" in slide_data and not has_bullet_points:
                # Use the same layout as bullet points would use
                left = Inches(0.5)
                top = Inches(1.8)
                width = Inches(left_column_width)
                height = Inches(5)
                
                # Create textbox
                textbox = slide.shapes.add_textbox(left, top, width, height)
                tf = textbox.text_frame
                configure_text_frame(tf)
                
                # Add paragraphs with rich formatting
                for para_data in slide_data["paragraphs"]:
                    add_text_with_mixed_formatting(tf, para_data)
                
                has_content_shape = True
            
            # Add chart if provided
            if has_chart:
                chart_data = slide_data["chart_data"]
                chart_type = chart_data.get("chart_type", "").upper() if isinstance(chart_data, dict) else ""
                
                # Position chart on right side if split, below content if stacked
                chart_left = Inches(right_column_left) if split_horizontally else Inches(0.5)
                chart_top = Inches(1.8) if split_horizontally else Inches(3.5)
                chart_width = Inches(right_column_width)
                chart_height = Inches(4.5) if split_horizontally else Inches(3.5)
                
                # Create appropriate chart type
                if chart_type in ["XY", "SCATTER"]:
                    xy_series = chart_data.get("xy_series", [])
                    if not xy_series and "series" in chart_data:
                        # Try to convert standard series to XY format
                        for series_name, values in chart_data["series"].items():
                            # Create points from values if possible
                            points = []
                            for i, y in enumerate(values):
                                points.append({"x": i, "y": y})
                            
                            xy_series.append({
                                "name": series_name,
                                "points": points
                            })
                    
                    chart = create_xy_chart(
                        slide,
                        x=chart_left.inches, 
                        y=chart_top.inches,
                        width=chart_width.inches,
                        height=chart_height.inches,
                        chart_data=xy_series,
                        title=chart_data.get("title", ""),
                        has_legend=chart_data.get("has_legend", True),
                        legend_position=chart_data.get("legend_position", "RIGHT"),
                        style=chart_data.get("style")
                    )
                elif chart_type == "BUBBLE":
                    bubble_series = chart_data.get("bubble_series", [])
                    
                    chart = create_bubble_chart(
                        slide,
                        x=chart_left.inches, 
                        y=chart_top.inches,
                        width=chart_width.inches,
                        height=chart_height.inches,
                        chart_data=bubble_series,
                        title=chart_data.get("title", ""),
                        has_legend=chart_data.get("has_legend", True),
                        legend_position=chart_data.get("legend_position", "RIGHT"),
                        style=chart_data.get("style")
                    )
                else:
                    # Standard category chart
                    categories = chart_data.get("categories", [])
                    series_data = chart_data.get("series", {})
                    
                    # Create chart data
                    cd = CategoryChartData()
                    cd.categories = categories
                    
                    # Add series
                    for series_name, values in series_data.items():
                        cd.add_series(series_name, values)
                    
                    # Map chart type string to PowerPoint chart type
                    chart_type_dict = {
                        "BAR": XL_CHART_TYPE.COLUMN_CLUSTERED,
                        "COLUMN": XL_CHART_TYPE.COLUMN_CLUSTERED,
                        "LINE": XL_CHART_TYPE.LINE,
                        "PIE": XL_CHART_TYPE.PIE,
                        "AREA": XL_CHART_TYPE.AREA,
                    }
                    xl_chart_type = chart_type_dict.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
                    
                    # Create chart
                    chart = slide.shapes.add_chart(
                        xl_chart_type, chart_left, chart_top, chart_width, chart_height, cd
                    ).chart
                    
                    # Set chart title if provided
                    chart.has_title = True
                    chart.chart_title.text_frame.text = chart_data.get("title", "")
                    
                    # Configure legend
                    chart.has_legend = chart_data.get("has_legend", True)
                    
                    # Set legend position if specified
                    legend_position = chart_data.get("legend_position", "RIGHT")
                    legend_pos_map = {
                        "right": XL_LEGEND_POSITION.RIGHT,
                        "left": XL_LEGEND_POSITION.LEFT,
                        "top": XL_LEGEND_POSITION.TOP,
                        "bottom": XL_LEGEND_POSITION.BOTTOM,
                        "corner": XL_LEGEND_POSITION.CORNER
                    }
                    if isinstance(legend_position, str) and legend_position.lower() in legend_pos_map:
                        chart.legend.position = legend_pos_map[legend_position.lower()]
                
                has_content_shape = True
            
            # Add table if provided
            if has_table:
                table_data = slide_data["table_data"]
                
                # Position table appropriately
                table_left = Inches(right_column_left) if (split_horizontally and has_chart) else Inches(0.5)
                table_top = Inches(1.8) if split_horizontally else (
                    Inches(3.5) if has_content_shape else Inches(1.8)
                )
                table_width = Inches(right_column_width if split_horizontally else 9.0)
                
                # Extract table components
                headers = table_data.get("headers", [])
                rows = table_data.get("rows", [])
                
                if headers and rows:
                    # Calculate reasonable height
                    num_rows = len(rows) + 1  # +1 for header
                    row_height = min(0.4, 4.0 / num_rows)
                    table_height = Inches(row_height * num_rows)
                    
                    # Create table
                    table = slide.shapes.add_table(
                        num_rows, len(headers), 
                        table_left, table_top, table_width, table_height
                    ).table
                    
                    # Add headers with formatting
                    for i, header in enumerate(headers):
                        cell = table.cell(0, i)
                        cell.text = str(header)
                        
                        # Format header text
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(0, 112, 192)  # Blue
                        
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.bold = True
                            paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White
                            paragraph.alignment = PP_ALIGN.CENTER
                    
                    # Add rows with formatting
                    for i, row_data in enumerate(rows):
                        for j, cell_text in enumerate(row_data):
                            if j < len(headers):  # Stay within bounds
                                cell = table.cell(i + 1, j)  # +1 to skip header row
                                cell.text = str(cell_text)
                                
                                # Format cell text
                                for paragraph in cell.text_frame.paragraphs:
                                    paragraph.font.size = Pt(11)
                                    paragraph.alignment = PP_ALIGN.CENTER
                                
                                # Apply alternating row colors
                                if i % 2 == 1:  # Every other row
                                    cell.fill.solid()
                                    cell.fill.fore_color.rgb = RGBColor(242, 242, 242)  # Light gray
                    
                    has_content_shape = True
            
            # Add image if provided
            if has_image:
                image_path = slide_data["image_path"]
                
                # Position image appropriately
                image_left = Inches(right_column_left) if (split_horizontally and (has_chart or has_table)) else Inches(2)
                image_top = Inches(1.8) if split_horizontally else (
                    Inches(3.5) if has_content_shape else Inches(1.8)
                )
                image_width = Inches(right_column_width if split_horizontally else 5.0)
                
                # Add image to slide
                slide.shapes.add_picture(image_path, image_left, image_top, width=image_width)
                
                has_content_shape = True
            
            # Add notes if provided
            if "notes" in slide_data and hasattr(slide, "notes_slide"):
                slide.notes_slide.notes_text_frame.text = slide_data["notes"]
                
            return slide
        except Exception as e:
            logging.warning(f"Error adding mixed content slide: {e}")
            logging.error(traceback.format_exc())
            return None

    def _add_connectors(self, prs: Presentation, connectors_data: List[Dict], shape_registry: Dict) -> None:
        """Add connectors between shapes.
        
        Args:
            prs: Presentation object
            connectors_data: List of connector configurations
            shape_registry: Dictionary of shapes indexed by ID
        """
        for connector_data in connectors_data:
            try:
                # Get required fields
                connector_type = connector_data.get("connector_type", "straight")
                begin_point = connector_data.get("begin_point")
                end_point = connector_data.get("end_point")
                
                # Get slide to add connector to
                slide = None
                
                # Check if we have shape connections
                begin_shape_id = connector_data.get("begin_shape_id")
                end_shape_id = connector_data.get("end_shape_id")
                
                if begin_shape_id and begin_shape_id in shape_registry:
                    slide = shape_registry[begin_shape_id].get("slide")
                elif end_shape_id and end_shape_id in shape_registry:
                    slide = shape_registry[end_shape_id].get("slide")
                
                # If we couldn't determine the slide, use the first slide
                if not slide and prs.slides:
                    slide = prs.slides[0]
                
                if slide and begin_point and end_point:
                    # Create the connector
                    connector = add_connector(
                        slide,
                        connector_type,
                        begin_point,
                        end_point,
                        line_color=connector_data.get("line_color", (0, 0, 0)),
                        line_width=connector_data.get("line_width", 1.5),
                        line_dash=connector_data.get("line_dash", "solid")
                    )
                    
                    # Add text if provided
                    if "text" in connector_data:
                        connector.text = connector_data["text"]
                        
                        # Configure text formatting
                        if "text_align" in connector_data:
                            align_map = {
                                "left": PP_ALIGN.LEFT,
                                "center": PP_ALIGN.CENTER,
                                "right": PP_ALIGN.RIGHT
                            }
                            if connector_data["text_align"].lower() in align_map:
                                connector.text_frame.paragraphs[0].alignment = align_map[connector_data["text_align"].lower()]
                        
                        if "font_bold" in connector_data:
                            connector.text_frame.paragraphs[0].font.bold = connector_data["font_bold"]
                        
                        if "font_size" in connector_data:
                            connector.text_frame.paragraphs[0].font.size = Pt(connector_data["font_size"])
                        
                        if "text_color" in connector_data:
                            set_font_color(connector.text_frame.paragraphs[0].font, connector_data["text_color"])
            except Exception as e:
                logging.warning(f"Error adding connector: {e}")

    def _save_presentation(self, prs: Presentation, title: str = "Presentation") -> Dict[str, str]:
        """Save the presentation to disk.
        
        Args:
            prs: Presentation object
            title: Presentation title for filename
            
        Returns:
            Dict with file paths
        """
        try:
            # Ensure output directory exists
            Path(self.output_dir).mkdir(parents=True, exist_ok=True)
            
            # Create unique filename
            safe_title = title.replace(" ", "_").replace(":", "")
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            unique_id = str(uuid.uuid4())[:8]
            filename = f"{safe_title}_{timestamp}_{unique_id}.pptx"
            filepath = os.path.join(self.output_dir, filename)
            abs_path = os.path.abspath(filepath)
                
            # Save the presentation
            prs.save(filepath)
            
            return {
                "file_path": abs_path,
                "relative_path": filepath
            }
        except Exception as e:
            logging.error(f"Error saving presentation: {e}")
            logging.error(traceback.format_exc())
            raise


class PythonPPTXTool(BaseTool):
    name: str = "PythonPPTXTool"
    description: str = (
        "A powerful tool for creating professional PowerPoint presentations using python-pptx. "
        "Features include:\n"
        "- Rich text formatting with mixed styles, colors, and fonts\n"
        "- Multiple chart types (bar, line, pie, XY/scatter, bubble)\n"
        "- Data tables with custom formatting and styling\n"
        "- Custom shapes with fill effects, gradients, and shadows\n"
        "- Connectors to link shapes with arrows\n"
        "- Hyperlinks and actions\n"
        "- Mixed content slides with multiple elements\n"
        "- Speaker notes for presentations\n\n"
        "Input can be plain text or structured JSON following the documented schema."
    )
    args_schema: Type[BaseModel] = PythonPPTXInput
    
    def __init__(self, **kwargs):
        super().__init__(**kwargs)
    
    def _run(self, content: str, title: Optional[str] = None,
             output_dir: str = "./", template_path: Optional[str] = None) -> PythonPPTXToolOutput:
        """Run the Python PPTX tool to create a presentation.
        
        Args:
            content: The raw content to be converted into slides. Can be plain text or JSON format.
            title: The title of the presentation
            output_dir: Directory to save the presentation (will be overridden to use ./output)
            template_path: Path to a PPTX template file (optional)
            
        Returns:
            PythonPPTXToolOutput: Output with file path and status
        """
        try:
            # Always use ./output directory regardless of what's provided
            output_dir = "./output"
            
            # Create the presentation generator
            generator = PPTXGenerator(
                output_dir=output_dir,
                template_path=template_path
            )
            
            # Parse and validate input
            actual_content = None
            
            if isinstance(content, str):
                try:
                    # Try to parse as JSON first
                    parsed_content = json.loads(content) if isinstance(content, str) else content
                    # Check if content has "articles" structure, which needs transformation
                    if "articles" in parsed_content:
                        # Transform articles format to PythonPPTX expected format
                        title_to_use = title or "Recent Events in Lebanon"
                        presentation_content = {
                            "title": title_to_use,
                            "headline": {
                                "title": title_to_use,
                                "subtitle": "News Summary",
                                "date": datetime.now().isoformat()
                            },
                            "slides": []
                        }
                        
                        # Create slides from articles
                        for article in parsed_content["articles"]:
                            slide = {
                                "title": article.get("title", "News Article"),
                                "content": article.get("description", ""),
                                "bullet_points": [
                                    {"text": "Source: " + article.get("company", "Unknown"), "bold": True},
                                ]
                            }
                            presentation_content["slides"].append(slide)
                        
                        actual_content = presentation_content
                    else:
                        # Check if the essential structure is present
                        if "title" not in parsed_content:
                            parsed_content["title"] = title or "Presentation"
                            
                        if "headline" not in parsed_content:
                            parsed_content["headline"] = {
                                "title": parsed_content["title"],
                                "subtitle": "Generated Presentation",
                                "date": datetime.now().isoformat()
                            }
                            
                        if "slides" not in parsed_content:
                            parsed_content["slides"] = [
                                {
                                    "title": "Overview",
                                    "content": "Generated content based on input"
                                }
                            ]
                        
                        actual_content = parsed_content
                except json.JSONDecodeError:
                    # Handle as simple text content - use title for headline
                    actual_title = title or "Presentation"
                    actual_content = {
                        "title": actual_title,
                        "headline": {
                            "title": actual_title,
                            "subtitle": "Generated Presentation",
                            "date": datetime.now().isoformat()
                        },
                        "slides": [
                            {
                                "title": "Content",
                                "content": content
                            }
                        ]
                    }
            else:
                # Already parsed content (dict)
                if isinstance(content, dict):
                    # Check if content has "articles" structure, which needs transformation
                    if "articles" in content:
                        # Transform articles format to PythonPPTX expected format
                        title_to_use = title or "Recent Events in Lebanon"
                        presentation_content = {
                            "title": title_to_use,
                            "headline": {
                                "title": title_to_use,
                                "subtitle": "News Summary",
                                "date": datetime.now().isoformat()
                            },
                            "slides": []
                        }
                        
                        # Create slides from articles
                        for article in content["articles"]:
                            slide = {
                                "title": article.get("title", "News Article"),
                                "content": article.get("description", ""),
                                "bullet_points": [
                                    {"text": "Source: " + article.get("company", "Unknown"), "bold": True},
                                ]
                            }
                            presentation_content["slides"].append(slide)
                        
                        actual_content = presentation_content
                    else:
                        # Check if the essential structure is present
                        if "title" not in content:
                            content["title"] = title or "Presentation"
                            
                        if "headline" not in content:
                            content["headline"] = {
                                "title": content["title"],
                                "subtitle": "Generated Presentation",
                                "date": datetime.now().isoformat()
                            }
                            
                        if "slides" not in content:
                            content["slides"] = [
                                {
                                    "title": "Overview",
                                    "content": "Generated content based on input"
                                }
                            ]
                        
                        actual_content = content
                else:
                    # Invalid content type
                    raise ValueError(f"Invalid content type: {type(content)}")
            
            # Generate the presentation using the transformed content
            result = generator.generate_from_json(actual_content, title)
            
            return PythonPPTXToolOutput(
                success=True,
                message="Presentation created successfully",
                file_path=result.get("file_path", ""),
                relative_path=result.get("relative_path", ""),
                content=json.dumps(actual_content) if isinstance(actual_content, dict) else content,
                title=actual_content.get("title", title) if isinstance(actual_content, dict) else title
            )
        except Exception as e:
            error_message = f"Error creating presentation: {str(e)}"
            logging.error(error_message)
            logging.error(traceback.format_exc())
            return PythonPPTXToolOutput(
                success=False,
                message=error_message,
                file_path="",
                relative_path="",
                content=content,
                title=title
            ) 