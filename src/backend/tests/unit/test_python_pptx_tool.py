import unittest
import os
import json
import tempfile
import shutil
from pathlib import Path
from datetime import datetime
from unittest.mock import patch, MagicMock
from pptx import Presentation

# Use relative imports that will work with the project structure
import sys
sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), '../../')))
from src.engines.crewai.tools.custom.python_pptx_tool import PythonPPTXTool
from src.engines.crewai.tools.schemas import PythonPPTXToolOutput

class TestPythonPPTXTool(unittest.TestCase):
    """Unit tests for PythonPPTXTool"""

    def setUp(self):
        """Set up test environment"""
        # Create a temporary directory for output
        self.test_dir = tempfile.mkdtemp()
        self.pptx_tool = PythonPPTXTool()
        
        # Sample presentation content
        self.sample_content = {
            "title": "Positive News from Lebanon",
            "author": "Lebanon News Network",
            "headline": {
                "title": "Lebanon Shows Signs of Economic Recovery",
                "subtitle": "New initiatives bring hope",
                "date": "2023-08-15T12:00:00"
            },
            "company": "Lebanese News & Media Group",
            "keywords": ["Lebanon", "economy", "culture", "tourism", "environment"],
            "slides": [
                {
                    "title": "Tourism Boom in Lebanon",
                    "bullet_points": [
                        {
                            "text": "Tourism increased by 30% in the past year",
                            "bold": True
                        },
                        {
                            "text": "Beirut named one of top cultural destinations by Travel Magazine",
                            "level": 1
                        },
                        {
                            "text": "Hotel occupancy rates reach 85% during summer months",
                            "level": 1
                        },
                        {
                            "text": "Revival of heritage sites attracting international visitors",
                            "italic": True
                        }
                    ],
                    "notes": "Emphasize the economic impact of tourism growth on local businesses"
                },
                {
                    "title": "Environmental Initiatives",
                    "bullet_points": [
                        "Cedar reforestation project plants 100,000 new trees",
                        "Beirut River cleanup removes 50 tons of waste",
                        "Solar power adoption increased 40% in rural areas",
                        "New marine protected areas established along coast"
                    ],
                    "notes": "Mention the international partnerships supporting these efforts"
                },
                {
                    "title": "Tourism Statistics",
                    "chart_data": {
                        "title": "Tourist Arrivals by Season (thousands)",
                        "chart_type": "BAR",
                        "categories": ["Winter", "Spring", "Summer", "Fall"],
                        "series": {
                            "2022": [120, 250, 480, 320],
                            "2023": [180, 310, 520, 390]
                        }
                    },
                    "notes": "Data courtesy of the Lebanon Tourism Board - highlight summer growth"
                },
                {
                    "title": "Economic Indicators",
                    "table_data": {
                        "headers": ["Indicator", "2021", "2022", "2023", "Change"],
                        "rows": [
                            ["GDP Growth", "-6.7%", "-2.1%", "1.8%", "↑"],
                            ["Exports (B $)", "3.2", "3.8", "4.5", "↑"],
                            ["Tourism Revenue (M $)", "420", "680", "950", "↑"],
                            ["Foreign Investment (M $)", "280", "520", "720", "↑"]
                        ]
                    },
                    "notes": "Source: Lebanon Central Bank and Ministry of Finance reports"
                },
                {
                    "title": "Cultural Achievements",
                    "bullet_points": [
                        "Lebanese film wins award at international festival",
                        "Beirut Art Fair attracts record number of visitors",
                        "Traditional Lebanese cuisine recognized by UNESCO",
                        "Lebanese designer showcased at Paris Fashion Week"
                    ],
                    "chart_data": {
                        "title": "Cultural Events Attendance",
                        "chart_type": "PIE",
                        "categories": ["Film Festivals", "Art Exhibitions", "Music Concerts", "Food Festivals"],
                        "series": {
                            "Attendees (thousands)": [45, 70, 120, 90]
                        }
                    }
                },
                {
                    "title": "Education Progress",
                    "bullet_points": [
                        "Literacy rate reaches 96% nationwide",
                        "New scholarship program supports 500 students",
                        "Lebanese universities establish international partnerships",
                        "Coding bootcamps graduate 1000+ new developers"
                    ],
                    "chart_data": {
                        "title": "Education Enrollment Trends",
                        "chart_type": "LINE",
                        "categories": ["2019", "2020", "2021", "2022", "2023"],
                        "series": {
                            "Higher Education": [55000, 56200, 58500, 61000, 65000],
                            "Technical Training": [12000, 14500, 18000, 22500, 28000]
                        }
                    }
                }
            ],
            "include_footer": True,
            "revision": 1
        }

    def tearDown(self):
        """Clean up after tests"""
        # Remove the temporary directory and its contents
        shutil.rmtree(self.test_dir)

    def test_run_with_json_content(self):
        """Test running the PPTX tool with JSON content"""
        # Convert sample content to JSON string
        json_content = json.dumps(self.sample_content)
        
        # Run the tool with the JSON content
        result = self.pptx_tool._run(
            content=json_content,
            title="Lebanon News Report",
            output_dir=self.test_dir
        )
        
        # Assert that the result is successful
        self.assertTrue(result.success)
        self.assertTrue(os.path.exists(result.file_path))
        
        # Verify the created presentation can be opened
        prs = Presentation(result.file_path)
        
        # Check basic structure of the presentation
        self.assertGreaterEqual(len(prs.slides), 7)  # Title slide + 6 content slides
        
        # Check title slide content
        title_slide = prs.slides[0]
        title_shapes = [shape for shape in title_slide.shapes if hasattr(shape, 'text') and shape.text]
        title_texts = [shape.text for shape in title_shapes]
        
        # At least one shape should contain the headline title
        self.assertTrue(any("Lebanon Shows Signs of Economic Recovery" in text for text in title_texts))

    def test_run_with_text_content_fallback(self):
        """Test running the PPTX tool with plain text content (fallback mode)"""
        text_content = """
        # Lebanon News Report
        
        Tourism is on the rise in Lebanon
        - Tourism increased by 30% last year
        - Beirut is becoming a popular destination
        
        # Environmental Projects
        
        Lebanon is improving its environment
        - Reforestation efforts underway
        - Cleanup initiatives showing success
        """
        
        # Run the tool with the text content
        result = self.pptx_tool._run(
            content=text_content,
            title="Lebanon Text Report",
            output_dir=self.test_dir
        )
        
        # Assert that the result is successful
        self.assertTrue(result.success)
        self.assertTrue(os.path.exists(result.file_path))
        
        # Verify the created presentation can be opened
        prs = Presentation(result.file_path)
        
        # Check if presentation has slides
        self.assertGreaterEqual(len(prs.slides), 1)

    def test_invalid_json_handling(self):
        """Test handling of invalid JSON content"""
        invalid_json = """
        {
            "title": "Invalid JSON Example",
            "content": "This JSON is missing a closing bracket,
            "slides": []
        """
        
        # Run the tool with invalid JSON
        result = self.pptx_tool._run(
            content=invalid_json,
            title="Invalid JSON Test",
            output_dir=self.test_dir
        )
        
        # Should still succeed using the fallback method
        self.assertTrue(result.success)
        self.assertTrue(os.path.exists(result.file_path))

    def test_exception_handling(self):
        """Test handling of exceptions during presentation creation"""
        # Create a mock tool that raises an exception
        with patch.object(PythonPPTXTool, '_run', side_effect=Exception("Forced test exception")):
            error_tool = PythonPPTXTool()
            
            try:
                # This should raise the exception we're forcing
                error_tool._run(
                    content=json.dumps(self.sample_content),
                    title="Exception Test",
                    output_dir=self.test_dir
                )
                self.fail("Expected an exception but none was raised")
            except Exception as e:
                self.assertIn("Forced test exception", str(e))
                
        # Alternative test with a real tool but mock an internal component
        with patch('src.engines.crewai.tools.custom.python_pptx_tool.Presentation', side_effect=Exception("Mock PowerPoint error")):
            result = self.pptx_tool._run(
                content=json.dumps(self.sample_content),
                title="Exception Test 2",
                output_dir=self.test_dir
            )
            
            # Should fail gracefully with error handling
            self.assertFalse(result.success)
            self.assertIn("error", result.message.lower())

if __name__ == '__main__':
    unittest.main() 